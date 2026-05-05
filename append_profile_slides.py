# -*- coding: utf-8 -*-
"""Append three profile-system slides to PF_Revit_Automation 1.pptx.

Matches the existing deck's design language verbatim:
  - 16:9 (13.33 x 7.5 in)
  - dark navy ground (set by slide master)
  - cyan accent  (#22D3EE) eyebrow bar + label, monospace Consolas
  - dark navy card  (#111B2E) with cyan left-edge stripe
  - Georgia titles (light  #E8EEF7)
  - Calibri body (muted #9CA9BF)
  - dark cyan numeric stamps (#0891B2)

The three new slides cover the profile system end-to-end:
  Slide A - Profile Capture     (one-and-done from a real fixture)
  Slide B - Profile Management  (YAML templates + ES storage)
  Slide C - Profile Placement   (before/after screenshots)
"""

import copy
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ---------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------

SRC_DIR = r"c:\CED_Extensions"
PPTX_IN = os.path.join(SRC_DIR, "PF_Revit_Automation 1.pptx")
PPTX_OUT = PPTX_IN  # overwrite in place; user asked to add to the same file

SCREENSHOT_DIR = (
    r"c:\Users\reed.pinterich\OneDrive - CoolSys Inc\Pictures\Screenshots"
)
IMG_BEFORE = os.path.join(SCREENSHOT_DIR, "Screenshot 2026-05-05 133540.png")  # 423x688
IMG_AFTER = os.path.join(SCREENSHOT_DIR, "Screenshot 2026-05-05 133509.png")   # 454x763


# ---------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------

COLOR_CYAN = RGBColor(0x22, 0xD3, 0xEE)     # accent / eyebrow / stripes
COLOR_CYAN_DARK = RGBColor(0x08, 0x91, 0xB2) # numeric stamps
COLOR_CARD = RGBColor(0x11, 0x1B, 0x2E)      # card fill
COLOR_LIGHT = RGBColor(0xE8, 0xEE, 0xF7)     # primary text
COLOR_MUTED = RGBColor(0x9C, 0xA9, 0xBF)     # secondary / body text


def _no_outline(shape):
    """Match the deck — none of the existing shapes carry a stroke."""
    try:
        shape.line.fill.background()
    except Exception:
        pass


def _add_filled_rect(slide, x_in, y_in, w_in, h_in, color):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x_in), Inches(y_in),
        Inches(w_in), Inches(h_in),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    _no_outline(sh)
    return sh


def _add_text(slide, x_in, y_in, w_in, h_in, text, *,
              font="Calibri", size_pt=11, bold=False, color=COLOR_MUTED,
              align=PP_ALIGN.LEFT):
    """Add a text-only auto shape (no fill, no stroke) — matches how
    the existing deck renders all of its labels."""
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x_in), Inches(y_in),
        Inches(w_in), Inches(h_in),
    )
    sh.fill.background()  # transparent
    _no_outline(sh)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return sh


def _add_eyebrow(slide, label):
    """The little cyan bar + ALL-CAPS Consolas label that opens every
    content slide in the deck."""
    _add_filled_rect(slide, 0.60, 0.50, 0.45, 0.06, COLOR_CYAN)
    _add_text(
        slide, 1.15, 0.32, 8.00, 0.40, label,
        font="Consolas", size_pt=11, bold=True, color=COLOR_CYAN,
    )


def _add_title(slide, text):
    """Big Georgia title under the eyebrow."""
    _add_text(
        slide, 0.60, 0.68, 12.00, 0.70, text,
        font="Georgia", size_pt=32, bold=True, color=COLOR_LIGHT,
    )


def _add_card(slide, x_in, y_in, w_in, h_in, *, number, title, body):
    """Reproduces the card pattern from the pipeline slide — cyan
    stripe, big number stamp top-right, title, body."""
    _add_filled_rect(slide, x_in, y_in, w_in, h_in, COLOR_CARD)
    _add_filled_rect(slide, x_in, y_in, 0.06, h_in, COLOR_CYAN)
    # Number stamp upper right.
    _add_text(
        slide, x_in + w_in - 1.10, y_in + 0.05, 1.00, 0.60, number,
        font="Georgia", size_pt=32, bold=True, color=COLOR_CYAN_DARK,
        align=PP_ALIGN.RIGHT,
    )
    # Title.
    _add_text(
        slide, x_in + 0.22, y_in + 0.18, w_in - 1.40, 0.45, title,
        font="Georgia", size_pt=15, bold=True, color=COLOR_LIGHT,
    )
    # Body.
    _add_text(
        slide, x_in + 0.22, y_in + 0.72, w_in - 0.40, h_in - 0.85, body,
        font="Calibri", size_pt=10.5, bold=False, color=COLOR_MUTED,
    )


def _add_image(slide, path, x_in, y_in, *, max_w_in=None, max_h_in=None):
    """Add an image, preserving its aspect ratio so 423x688 vs
    454x763 don't look stretched."""
    # python-pptx infers the natural aspect when only one dim is given.
    kwargs = {}
    if max_w_in is not None and max_h_in is not None:
        # Pick the smaller scale to fit both bounds.
        from PIL import Image
        with Image.open(path) as im:
            w_px, h_px = im.size
        scale = min(max_w_in / (w_px / 96.0), max_h_in / (h_px / 96.0))
        # Use 96 DPI as a safe anchor for screen-shot pixel-to-inch.
        kwargs["width"] = Inches((w_px / 96.0) * scale)
        kwargs["height"] = Inches((h_px / 96.0) * scale)
    elif max_w_in is not None:
        kwargs["width"] = Inches(max_w_in)
    elif max_h_in is not None:
        kwargs["height"] = Inches(max_h_in)
    return slide.shapes.add_picture(path, Inches(x_in), Inches(y_in), **kwargs)


# ---------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------

def build_slide_capture(slide):
    _add_eyebrow(slide, "PROFILE CAPTURE")
    _add_title(slide, "Capture once, place forever.")

    _add_text(
        slide, 0.60, 1.40, 12.00, 0.55,
        "Walk a real fixture in Revit; the panel records every linked "
        "child, parameter value, offset, and annotation as a YAML "
        "profile. The next project starts from that captured truth, "
        "not a blank slate.",
        font="Calibri", size_pt=14, bold=False, color=COLOR_MUTED,
    )

    # Three cards across, mirroring the pipeline layout (4.05 x 2.65).
    y = 2.20
    h = 2.65
    _add_card(
        slide, 0.60, y, 4.05, h,
        number="01",
        title="Pick the parent",
        body=(
            "User selects an equipment parent (case, freezer, "
            "millwork, etc.). The agent walks every linked element via "
            "GetDependentElements and explicit picks; nothing else is "
            "guessed.\n\n"
            "Element_Linker JSON gets stamped on every captured "
            "element so audits can prove it back to its profile later."
        ),
    )
    _add_card(
        slide, 4.80, y, 4.05, h,
        number="02",
        title="Encode offsets + params",
        body=(
            "Geometry is captured in the parent's local frame, so the "
            "whole assembly rotates as a rigid body when placed at a "
            "different orientation.\n\n"
            "Every parameter the user can edit on the live element is "
            "written into the YAML — including type-level params for "
            "keynote families (Keynote Value, Description)."
        ),
    )
    _add_card(
        slide, 9.00, y, 4.05, h,
        number="03",
        title="Tags, keynotes, text",
        body=(
            "Annotations travel with their host: family : type, world "
            "offset relative to the host fixture, leaders for tags, "
            "literal text content for text notes.\n\n"
            "On placement the annotation list rebuilds at the chosen "
            "location. Dedup compares against any tags / keynotes "
            "already on the new instance."
        ),
    )

    # Bottom stat band — same style as slide 1's chip row.
    _add_filled_rect(slide, 0.57, 5.85, 6.00, 1.15, COLOR_CARD)
    _add_filled_rect(slide, 0.57, 5.85, 0.06, 1.15, COLOR_CYAN)
    _add_text(
        slide, 0.78, 5.97, 5.70, 0.55, "30+ profiles",
        font="Georgia", size_pt=28, bold=True, color=COLOR_LIGHT,
    )
    _add_text(
        slide, 0.78, 6.53, 5.70, 0.40,
        "captured for HEB and Planet Fitness so far",
        font="Calibri", size_pt=11, bold=False, color=COLOR_MUTED,
    )

    _add_filled_rect(slide, 6.78, 5.85, 6.00, 1.15, COLOR_CARD)
    _add_filled_rect(slide, 6.78, 5.85, 0.06, 1.15, COLOR_CYAN)
    _add_text(
        slide, 6.97, 5.97, 5.70, 0.55, "0 hand re-modeling",
        font="Georgia", size_pt=28, bold=True, color=COLOR_LIGHT,
    )
    _add_text(
        slide, 6.97, 6.53, 5.70, 0.40,
        "every linked element + annotation is captured automatically",
        font="Calibri", size_pt=11, bold=False, color=COLOR_MUTED,
    )


def build_slide_management(slide):
    _add_eyebrow(slide, "PROFILE MANAGEMENT")
    _add_title(slide, "YAML in, YAML out.")

    _add_text(
        slide, 0.60, 1.40, 12.00, 0.55,
        "Profiles live in two places: a per-project Extensible "
        "Storage entity (so the data ships with the model) and an "
        "exportable YAML file (so the same templates seed every "
        "project that follows).",
        font="Calibri", size_pt=14, bold=False, color=COLOR_MUTED,
    )

    y = 2.20
    h = 2.65
    _add_card(
        slide, 0.60, y, 4.05, h,
        number="01",
        title="Manage Profiles editor",
        body=(
            "Tree view: linked sets, LEDs, annotations. The user "
            "edits parameter values, offsets, family : type, and "
            "merge-group truth-source pointers — every change "
            "round-trips cleanly through the YAML.\n\n"
            "Captured fields display + edit live; built-in keynote "
            "values surface alongside shared CKT_* params."
        ),
    )
    _add_card(
        slide, 4.80, y, 4.05, h,
        number="02",
        title="Extensible Storage v4",
        body=(
            "Project-local data lives on a dedicated DataStorage "
            "element (not ProjectInformation) so it survives schema "
            "changes and stays out of users' way.\n\n"
            "Four typed map fields (String / Bool / Int64 / Double) "
            "+ DocGuid. Int64-first with Int32 fallback keeps Revit "
            "2026 happy."
        ),
    )
    _add_card(
        slide, 9.00, y, 4.05, h,
        number="03",
        title="Import / export YAML",
        body=(
            "Output is byte-identical to what's stored — Import then "
            "Export is lossless. Same template can seed five "
            "projects with no copy/paste.\n\n"
            "Spaces config exports separately from equipment so a "
            "bakery template moves to the next store without "
            "carrying that store's space-to-bucket assignments."
        ),
    )

    _add_filled_rect(slide, 0.57, 5.85, 6.00, 1.15, COLOR_CARD)
    _add_filled_rect(slide, 0.57, 5.85, 0.06, 1.15, COLOR_CYAN)
    _add_text(
        slide, 0.78, 5.97, 5.70, 0.55, "1 source of truth",
        font="Georgia", size_pt=28, bold=True, color=COLOR_LIGHT,
    )
    _add_text(
        slide, 0.78, 6.53, 5.70, 0.40,
        "edits in the YAML propagate to every project on next import",
        font="Calibri", size_pt=11, bold=False, color=COLOR_MUTED,
    )

    _add_filled_rect(slide, 6.78, 5.85, 6.00, 1.15, COLOR_CARD)
    _add_filled_rect(slide, 6.78, 5.85, 0.06, 1.15, COLOR_CYAN)
    _add_text(
        slide, 6.97, 5.97, 5.70, 0.55, "v100 schema",
        font="Georgia", size_pt=28, bold=True, color=COLOR_LIGHT,
    )
    _add_text(
        slide, 6.97, 6.53, 5.70, 0.40,
        "with v3 / v4 legacy migration on import",
        font="Calibri", size_pt=11, bold=False, color=COLOR_MUTED,
    )


def build_slide_placement(slide):
    _add_eyebrow(slide, "PROFILE PLACEMENT")
    _add_title(slide, "From AutoCAD to Revit, in seconds.")

    _add_text(
        slide, 0.60, 1.40, 12.00, 0.55,
        "Place from CAD or a linked Revit model: the engine matches "
        "every parent target to a captured profile, drops every "
        "linked element at its captured offset, and stamps the "
        "Element_Linker JSON for downstream audits and circuiting.",
        font="Calibri", size_pt=14, bold=False, color=COLOR_MUTED,
    )

    # Two-up before/after panel layout.
    img_top = 2.15
    img_h = 4.50
    panel_w = 5.80
    panel_x_left = 0.60
    panel_x_right = 6.93

    # Backing cards for each image.
    _add_filled_rect(slide, panel_x_left, img_top - 0.05,
                     panel_w, img_h + 0.55, COLOR_CARD)
    _add_filled_rect(slide, panel_x_left, img_top - 0.05,
                     0.06, img_h + 0.55, COLOR_CYAN)
    _add_filled_rect(slide, panel_x_right, img_top - 0.05,
                     panel_w, img_h + 0.55, COLOR_CARD)
    _add_filled_rect(slide, panel_x_right, img_top - 0.05,
                     0.06, img_h + 0.55, COLOR_CYAN)

    # Labels: BEFORE / AFTER.
    _add_text(
        slide, panel_x_left + 0.20, img_top + img_h + 0.05,
        panel_w - 0.30, 0.30,
        "BEFORE  -  CAD layout (provided by architect)",
        font="Consolas", size_pt=11, bold=True, color=COLOR_CYAN,
    )
    _add_text(
        slide, panel_x_right + 0.20, img_top + img_h + 0.05,
        panel_w - 0.30, 0.30,
        "AFTER  -  Revit equipment placed automatically",
        font="Consolas", size_pt=11, bold=True, color=COLOR_CYAN,
    )

    # Images centered within their panels, fitted to the box while
    # preserving aspect.
    _fit_image(
        slide, IMG_BEFORE,
        panel_x_left + 0.20, img_top + 0.05,
        panel_w - 0.40, img_h - 0.10,
    )
    _fit_image(
        slide, IMG_AFTER,
        panel_x_right + 0.20, img_top + 0.05,
        panel_w - 0.40, img_h - 0.10,
    )


def _fit_image(slide, path, x_in, y_in, max_w_in, max_h_in):
    """Insert ``path`` centered in the ``max_w_in`` x ``max_h_in`` box
    starting at (x_in, y_in). Aspect ratio preserved from PNG pixels;
    the screenshot dpi assumption is 96 dpi (Windows default)."""
    from PIL import Image
    with Image.open(path) as im:
        w_px, h_px = im.size
    nat_w_in = w_px / 96.0
    nat_h_in = h_px / 96.0
    scale = min(max_w_in / nat_w_in, max_h_in / nat_h_in)
    img_w_in = nat_w_in * scale
    img_h_in = nat_h_in * scale
    cx = x_in + (max_w_in - img_w_in) / 2.0
    cy = y_in + (max_h_in - img_h_in) / 2.0
    slide.shapes.add_picture(
        path, Inches(cx), Inches(cy),
        width=Inches(img_w_in), height=Inches(img_h_in),
    )


# ---------------------------------------------------------------------
# Append
# ---------------------------------------------------------------------

def main():
    prs = Presentation(PPTX_IN)
    blank_layout = prs.slide_layouts[0]  # only one layout in the deck

    builders = [
        build_slide_capture,
        build_slide_management,
        build_slide_placement,
    ]
    for build in builders:
        slide = prs.slides.add_slide(blank_layout)
        build(slide)

    prs.save(PPTX_OUT)
    print("Wrote {} ({} slides total)".format(PPTX_OUT, len(prs.slides)))


if __name__ == "__main__":
    main()
