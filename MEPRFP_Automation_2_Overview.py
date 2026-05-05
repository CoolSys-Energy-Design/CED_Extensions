# -*- coding: utf-8 -*-
"""Generate a brief MEPRFP Automation 2.0 overview deck.

Runs in plain CPython (not under pyRevit) — uses python-pptx to write
``MEPRFP_Automation_2_Overview.pptx`` next to this file.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ---------------------------------------------------------------------
# Slide content (title, body lines)
# ---------------------------------------------------------------------

_BRAND_BLUE = RGBColor(0x1F, 0x3A, 0x5F)
_ACCENT = RGBColor(0xC0, 0x39, 0x2B)


def _slides():
    return [
        {
            "title": "MEPRFP Automation 2.0",
            "subtitle": "A Revit panel for repeatable MEP layouts\n"
                        "Equipment profiles · Spaces · Circuiting",
            "kind": "title",
        },
        {
            "title": "Why we built this",
            "bullets": [
                "MEP layouts are repetitive across stores / project types — "
                "the same fixtures land in the same relative spots over and over.",
                "Hand-laying each instance burns hours and produces inconsistent "
                "output across projects and engineers.",
                "MEPRFP Automation 2.0 lets us capture a layout once as a "
                "reusable template, then place it everywhere it's needed.",
                "Templates live as project-local data plus exportable YAML, so "
                "they round-trip across projects cleanly.",
            ],
        },
        {
            "title": "What's in the panel",
            "bullets": [
                "Import-Export Profiles — YAML in / out for sharing.",
                "Modify Profiles — author and edit equipment templates.",
                "Place from CAD or Linked Model — drop equipment from a "
                "host model or DWG block layout.",
                "Place Single Profile — one-off placement at a picked point.",
                "Place Element Annotations — tags, keynotes, text notes.",
                "Spaces — bucket-based templates anchored to room geometry.",
                "Circuiting — SuperCircuit V5 builds Revit electrical systems "
                "from the captured fixtures.",
                "Audit tools — find drift between YAML and the live model.",
            ],
        },
        {
            "title": "The profile model",
            "bullets": [
                "A profile is a YAML record describing a parent fixture and "
                "the linked elements (LEDs) that travel with it.",
                "Each LED carries: family : type label, parameters, offsets, "
                "and an annotation list (tags / keynotes / text notes).",
                "Offsets are expressed in the parent's local frame — so when "
                "a parent rotates, every linked element rotates with it.",
                "Profiles support truth-source merge groups so a fix in one "
                "definition propagates to every member of its group.",
            ],
        },
        {
            "title": "Where the data lives",
            "bullets": [
                "All project-local data sits in Extensible Storage on a "
                "dedicated DataStorage element (not on ProjectInformation).",
                "Two ES schemas — equipment YAML and space classifications — "
                "are kept on separate DataStorages so an export of one "
                "doesn't drag the other.",
                "v4 schema layout: four typed map fields (String / Bool / "
                "Int64 / Double) plus a DocGuid string. Int64 with Int32 "
                "fallback keeps Revit 2026 happy.",
                "Legacy v1 entities (simple-fields layout on "
                "ProjectInformation) still read; first save migrates the "
                "data forward without GUID changes.",
            ],
        },
        {
            "title": "Equipment workflow — author",
            "bullets": [
                "Select a parent fixture in Revit + run Capture.",
                "The capture engine walks every linked element, records "
                "family : type, parameters, offsets relative to the parent, "
                "and any tags / keynotes / text notes near it.",
                "The captured profile lands in the active project's YAML "
                "store — Manage Profiles lets you rename, merge, edit, and "
                "delete entries afterward.",
                "Element_Linker JSON is stamped on every captured element "
                "so audits can match each placed element back to its profile.",
            ],
        },
        {
            "title": "Equipment workflow — place",
            "bullets": [
                "Place from CAD or Linked Model: pick a host doc (linked "
                "Revit, DWG block layout, or CSV) and the engine matches "
                "every parent target to a profile, then drops every LED at "
                "its captured offset.",
                "Place Single Profile: pick a point in the host view, "
                "choose a profile from the active store, and the engine "
                "creates the family instances at the offsets.",
                "Place Element Annotations: walks placed fixtures, looks up "
                "their LED's annotations list, and creates tags / keynotes / "
                "text notes — already-placed items are deduped by default.",
            ],
        },
        {
            "title": "Circuiting — SuperCircuit V5",
            "bullets": [
                "Reads the CKT_Panel_CEDT, CKT_Circuit Number_CEDT, "
                "CKT_Load Name_CEDT, CKT_Rating_CED, and CKT_Schedule "
                "Notes_CEDT parameters off captured fixtures.",
                "Groups fixtures into circuits (DEDICATED, BYPARENT, "
                "SECONDBYPARENT, NORMAL) and creates Revit "
                "ElectricalSystem objects via an ExternalEvent gateway.",
                "Client-aware: HEB and Planet Fitness keyword sets are "
                "handled by per-client classifier modules with a shared base.",
                "Audit Circuits surfaces drift between YAML and live "
                "circuits — missing data, phantom panels, orphan circuits, "
                "pole mismatches.",
            ],
        },
        {
            "title": "Spaces — bucket vocabulary",
            "bullets": [
                "Buckets are keyword-tagged categories like RESTROOM, "
                "BAKERY, ELECTRICAL ROOM. Each bucket has a list of "
                "case-insensitive substring keywords + an optional "
                "client-key restriction.",
                "Manage Space Buckets lets the user define / edit the "
                "vocabulary directly in Revit.",
                "Classify Spaces walks every placed Revit Space, matches "
                "its name against bucket keywords, and lets the user "
                "override per-space (multi-bucket assignment supported).",
                "Classifications are project-local — they don't ride with "
                "exported templates.",
            ],
        },
        {
            "title": "Spaces — profiles & door-aware anchors",
            "bullets": [
                "A space profile is a template entry that targets a bucket "
                "and lists LEDs to place inside any matching space.",
                "Multiple profiles can target the same bucket — they stack, "
                "so a Restroom space gets every Restroom profile's elements.",
                "Anchor kinds (no cardinals — door-aware): center, "
                "door_relative, wall_opposite_door, wall_right_of_door, "
                "wall_left_of_door, corner_furthest_from_door, "
                "corner_closest_to_door.",
                "Manage Space Profiles edits buckets, profiles, LEDs, "
                "parameters, offsets, and annotations end-to-end without "
                "touching YAML by hand.",
            ],
        },
        {
            "title": "Spaces — placement flow",
            "bullets": [
                "Place Space Elements collects every classified space, "
                "finds matching profiles, and previews every prospective "
                "family instance in a flat selectable table.",
                "If a space has more than one door AND a door-dependent "
                "anchor, the user is prompted to click the reference door "
                "in the model (host or linked) — that choice is reused for "
                "every door-anchored LED in that space.",
                "Per-row Select checkbox + Bucket column let the user place "
                "subsets (e.g. all bakery profiles, none of the freezer "
                "ones) without re-classifying.",
                "Place Space Annotations walks placed space-based fixtures "
                "and drops their tags / keynotes — same dedup machinery "
                "as the equipment side.",
            ],
        },
        {
            "title": "Sharing across projects",
            "bullets": [
                "Import / Export YAML File round-trips equipment definitions "
                "between projects. Output is byte-identical to the stored "
                "payload — Import → Export is lossless.",
                "Import / Export Space Config writes only the space-template "
                "data (buckets + profiles), never the per-project "
                "classifications. A bakery template moves to the next store "
                "without taking that store's space-to-bucket assignments "
                "with it.",
                "Blank-file imports synthesise an empty v100 payload, so a "
                "fresh project can start empty and be built up using the "
                "in-Revit editors.",
            ],
        },
        {
            "title": "Element_Linker — the audit backbone",
            "bullets": [
                "Every placed element gets a JSON payload written to its "
                "Element_Linker shared parameter at placement time.",
                "Equipment fields: led_id, set_id, location, rotation, "
                "parent rotation, host name, level, facing, parent location, "
                "ckt_panel, ckt_circuit_number.",
                "Spaces fields (Stage 6): space_id, space_profile_id — "
                "lets audits prove every space-based element back to the "
                "Space and template that produced it.",
                "Audit / sync tools use this lineage to detect drift, "
                "phantom elements, and missing-data conditions.",
            ],
        },
        {
            "title": "What's next",
            "bullets": [
                "Smoke-test the full Stage 6 + Stage 7 flow in a real "
                "project (in progress).",
                "Optional bucket auto-rerun on classification changes — "
                "today re-classification is manual.",
                "Field-feedback iteration on the door picker UX (status-bar "
                "prompts vs. a persistent overlay).",
                "Extend audits to surface profile-vs-element drift on the "
                "Spaces side (parity with equipment-side Sync Audit).",
                "Continue documenting per-client keyword sets so new "
                "clients onboard via a small adapter rather than core code.",
            ],
        },
    ]


# ---------------------------------------------------------------------
# Builder
# ---------------------------------------------------------------------

def _set_title_text(title_shape, text, color=_BRAND_BLUE, size=32):
    tf = title_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color


def _add_bullets(body_shape, bullets, font_size=16):
    tf = body_shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)


def _add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]  # title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    title_run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(40)
    title_run.font.bold = True
    title_run.font.color.rgb = _BRAND_BLUE

    if len(slide.placeholders) > 1:
        sub = slide.placeholders[1]
        sub.text = subtitle
        for p in sub.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(18)
                r.font.color.rgb = _ACCENT


def _add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]  # title + content
    slide = prs.slides.add_slide(slide_layout)
    _set_title_text(slide.shapes.title, title)
    body = slide.placeholders[1]
    _add_bullets(body, bullets)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for s in _slides():
        if s.get("kind") == "title":
            _add_title_slide(prs, s["title"], s["subtitle"])
        else:
            _add_content_slide(prs, s["title"], s["bullets"])

    out_path = "MEPRFP_Automation_2_Overview.pptx"
    prs.save(out_path)
    print("Wrote {} ({} slides)".format(out_path, len(prs.slides)))


if __name__ == "__main__":
    main()
